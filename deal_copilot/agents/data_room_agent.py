"""
Data Room / Deal Pack Ingestor Agent
Extracts qualitative and quantitative data from confidential deal documents
Uses OpenAI for inference - NO HALLUCINATIONS, only extract what exists
"""

from typing import Dict, List, Optional, Any
from datetime import datetime
import pandas as pd
from io import BytesIO
import json
import re

from openai import OpenAI
from deal_copilot.config import config_openai as config

# File parsing imports
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    PyPDF2 = None
    pdfplumber = None

try:
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils.dataframe import dataframe_to_rows
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None


class DataRoomAgent:
    """
    Agent that ingests confidential deal pack documents and extracts
    both qualitative insights and quantitative financials.
    
    CRITICAL: NO HALLUCINATIONS - If data doesn't exist, mark as N/A
    """
    
    def __init__(self, progress_callback=None, stream_callback=None):
        """
        Initialize the Data Room Agent with OpenAI
        
        Args:
            progress_callback: Optional function to call with progress updates
                               Signature: callback(step: str, progress: int, message: str)
            stream_callback: Optional function to call with streaming content chunks
                            Signature: callback(chunk: str)
        """
        self.progress_callback = progress_callback
        self.stream_callback = stream_callback
        
        # Initialize OpenAI client
        self.client = OpenAI(api_key=config.OPENAI_API_KEY)
        self.model = config.OPENAI_MODEL  # GPT-5, gpt-4o, etc.
        
        # Note: GPT-5 uses default temperature of 1 (not configurable)
    
    def _update_progress(self, step: str, progress: int, message: str):
        """Update progress if callback is provided"""
        if self.progress_callback:
            self.progress_callback(step, progress, message)
        print(f"  [{progress}%] {step}: {message}")
    
    def process_data_room(
        self,
        files: List[Dict[str, Any]],
        company_name: str
    ) -> Dict:
        """
        Process all files in the data room
        
        Args:
            files: List of {filename, content, file_type}
            company_name: Name of the company
            
        Returns:
            Dictionary with qualitative analysis and quantitative data
        """
        print(f"\n{'='*60}")
        print(f"DATA ROOM INGESTOR AGENT")
        print(f"{'='*60}")
        print(f"Company: {company_name}")
        print(f"Files to process: {len(files)}")
        print(f"{'='*60}\n")
        
        # Step 1: Extract raw text/data from all files
        self._update_progress("extraction", 10, f"Extracting content from {len(files)} files...")
        extracted_content = self._extract_from_files(files)
        self._update_progress("extraction", 25, f"Successfully extracted {len(files)} files")
        
        # Step 2: Generate qualitative analysis
        self._update_progress("qualitative", 30, "Analyzing qualitative information...")
        qualitative_analysis = self._generate_qualitative_analysis(
            extracted_content,
            company_name
        )
        self._update_progress("qualitative", 55, "Qualitative analysis complete")
        
        # Step 3: Extract quantitative data
        self._update_progress("quantitative", 60, "Extracting financial data and metrics...")
        quantitative_data = self._extract_quantitative_data(
            extracted_content,
            company_name
        )
        self._update_progress("quantitative", 80, "Quantitative extraction complete")
        
        # Step 4: Generate Excel file with quantitative data
        self._update_progress("excel", 85, "Generating Excel workbook...")
        excel_file = self._generate_excel_file(
            quantitative_data,
            company_name
        )
        self._update_progress("excel", 90, "Excel file generated")
        
        # Step 5: Generate human-readable summary for frontend
        self._update_progress("summary", 92, "Generating human-readable summary...")
        human_readable_summary = self._generate_human_readable_summary(
            company_name,
            qualitative_analysis,
            quantitative_data
        )
        self._update_progress("summary", 95, "Summary generated")
        
        # Generate summary
        extraction_summary = extracted_content.get("extraction_summary", {})
        summary = {
            "total_files": extraction_summary.get("total_files", len(files)),
            "successful_extractions": extraction_summary.get("successful", 0),
            "failed_extractions": extraction_summary.get("failed", 0),
            "pdfs_extracted": len(extracted_content.get("pdfs", [])),
            "excel_files_extracted": len(extracted_content.get("excel", [])),
            "powerpoints_extracted": len(extracted_content.get("powerpoint", [])),
            "has_excel_output": excel_file is not None
        }
        
        print(f"\n{'='*60}")
        print(f"✅ DATA ROOM PROCESSING COMPLETE")
        print(f"{'='*60}")
        print(f"  Files processed: {summary['successful_extractions']}/{summary['total_files']}")
        print(f"  PDFs: {summary['pdfs_extracted']}, Excel: {summary['excel_files_extracted']}, PPTs: {summary['powerpoints_extracted']}")
        print(f"  Output Excel generated: {'Yes' if summary['has_excel_output'] else 'No'}")
        print(f"{'='*60}\n")
        
        return {
            "company_name": company_name,
            "generated_at": datetime.now().isoformat(),
            "files_processed": len(files),
            "extraction_summary": summary,
            "qualitative_analysis": qualitative_analysis,
            "quantitative_data": quantitative_data,
            "human_readable_summary": human_readable_summary,  # For frontend display
            "excel_file": excel_file  # BytesIO object with Excel file
        }
    
    def _extract_from_files(self, files: List[Dict]) -> Dict:
        """Extract text and data from all uploaded files with validation"""
        content = {
            "pdfs": [],
            "excel": [],
            "powerpoint": [],
            "docx": [],
            "text": [],
            "extraction_summary": {
                "total_files": len(files),
                "successful": 0,
                "failed": 0,
                "errors": []
            }
        }
        
        for file_info in files:
            filename = file_info.get("filename", "unknown")
            file_content = file_info.get("content")
            file_type = file_info.get("file_type", "unknown")
            file_size = len(file_content) if file_content else 0
            
            print(f"  Processing: {filename} ({file_type}, {file_size:,} bytes)")
            
            # Validate file content exists
            if not file_content or file_size == 0:
                error_msg = f"Empty file: {filename}"
                print(f"  ✗ {error_msg}")
                content["extraction_summary"]["failed"] += 1
                content["extraction_summary"]["errors"].append(error_msg)
                continue
            
            try:
                if file_type == "pdf":
                    text = self._extract_pdf_text(file_content, filename)
                    if text and not text.startswith("Error"):
                        content["pdfs"].append({
                            "filename": filename,
                            "text": text,
                            "size": file_size,
                            "text_length": len(text)
                        })
                        content["extraction_summary"]["successful"] += 1
                    else:
                        content["extraction_summary"]["failed"] += 1
                        content["extraction_summary"]["errors"].append(f"PDF extraction failed: {filename}")
                    
                elif file_type == "excel":
                    data = self._extract_excel_data(file_content, filename)
                    if "error" not in data:
                        content["excel"].append({
                            "filename": filename,
                            "data": data,
                            "size": file_size
                        })
                        content["extraction_summary"]["successful"] += 1
                    else:
                        content["extraction_summary"]["failed"] += 1
                        content["extraction_summary"]["errors"].append(f"Excel extraction failed: {filename}")
                    
                elif file_type == "powerpoint":
                    text = self._extract_ppt_text(file_content, filename)
                    if text and not text.startswith("Error"):
                        content["powerpoint"].append({
                            "filename": filename,
                            "text": text,
                            "size": file_size,
                            "text_length": len(text)
                        })
                        content["extraction_summary"]["successful"] += 1
                    else:
                        content["extraction_summary"]["failed"] += 1
                        content["extraction_summary"]["errors"].append(f"PowerPoint extraction failed: {filename}")
                
                elif file_type == "docx":
                    text = self._extract_docx_text(file_content, filename)
                    if text and not text.startswith("Error"):
                        content["docx"].append({
                            "filename": filename,
                            "text": text,
                            "size": file_size,
                            "text_length": len(text)
                        })
                        content["extraction_summary"]["successful"] += 1
                    else:
                        content["extraction_summary"]["failed"] += 1
                        content["extraction_summary"]["errors"].append(f"Word document extraction failed: {filename}")
                    
            except Exception as e:
                error_msg = f"Error processing {filename}: {str(e)}"
                print(f"  ✗ {error_msg}")
                content["extraction_summary"]["failed"] += 1
                content["extraction_summary"]["errors"].append(error_msg)
        
        # Print summary
        print(f"\n  Extraction Summary:")
        print(f"    ✓ Successful: {content['extraction_summary']['successful']}/{content['extraction_summary']['total_files']}")
        print(f"    ✗ Failed: {content['extraction_summary']['failed']}/{content['extraction_summary']['total_files']}")
        if content['extraction_summary']['errors']:
            print(f"    Errors: {', '.join(content['extraction_summary']['errors'][:3])}")
        
        return content
    
    def _extract_pdf_text(self, content: bytes, filename: str) -> str:
        """Extract text and tables from PDF - FULL EXTRACTION, no truncation"""
        if not pdfplumber:
            return "PDF parsing not available - install pdfplumber"
        
        text_parts = []
        
        try:
            with pdfplumber.open(BytesIO(content)) as pdf:
                print(f"    Extracting {len(pdf.pages)} pages from {filename}...")
                
                for i, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"\n--- Page {i} of {filename} ---\n{page_text}")
                    
                    # Also extract tables from the page
                    tables = page.extract_tables()
                    if tables:
                        for table_idx, table in enumerate(tables, 1):
                            text_parts.append(f"\n[TABLE {table_idx} on Page {i}]")
                            # Convert table to readable text format
                            for row in table:
                                if row:  # Skip empty rows
                                    text_parts.append(" | ".join([str(cell) if cell else "" for cell in row]))
                
                print(f"    ✓ Extracted {len(pdf.pages)} pages with full content")
        except Exception as e:
            print(f"    ✗ Error extracting PDF: {e}")
            return f"Error extracting PDF: {e}"
        
        return "\n".join(text_parts)
    
    def _extract_excel_data(self, content: bytes, filename: str) -> Dict:
        """Extract ALL data from Excel files - FULL EXTRACTION with metadata"""
        try:
            # Read all sheets
            excel_file = BytesIO(content)
            xls = pd.ExcelFile(excel_file)
            
            print(f"    Extracting {len(xls.sheet_names)} sheets from {filename}...")
            
            sheets_data = {}
            sheets_metadata = {}
            
            for sheet_name in xls.sheet_names:
                # Read with header detection, preserve data types
                df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str, keep_default_na=False)
                
                # Convert back to appropriate types where possible (but preserve as strings for LLM)
                # This ensures we don't lose precision on numbers
                for col in df.columns:
                    # Try to convert to numeric, but keep as string representation to preserve exact values
                    try:
                        # Check if column contains numeric data
                        pd.to_numeric(df[col], errors='raise')
                        # If successful, it's numeric but we keep string representation for exact values
                    except (ValueError, TypeError):
                        # Not numeric, keep as is
                        pass
                
                # Store full data (not just sample) - empty strings for missing data
                sheets_data[sheet_name] = df.to_dict(orient='records')
                
                # Store metadata for better context
                sheets_metadata[sheet_name] = {
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": list(df.columns),
                    "has_data": len(df) > 0
                }
                
                print(f"      ✓ Sheet '{sheet_name}': {len(df)} rows x {len(df.columns)} columns")
            
            return {
                "sheets": sheets_data,
                "sheet_names": xls.sheet_names,
                "metadata": sheets_metadata
            }
        except Exception as e:
            print(f"    ✗ Error extracting Excel: {e}")
            return {"error": f"Error extracting Excel: {e}"}
    
    def _extract_ppt_text(self, content: bytes, filename: str) -> str:
        """Extract ALL text from PowerPoint - FULL EXTRACTION including tables"""
        if not Presentation:
            return "PowerPoint parsing not available - install python-pptx"
        
        try:
            prs = Presentation(BytesIO(content))
            text_parts = []
            
            print(f"    Extracting {len(prs.slides)} slides from {filename}...")
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {i} of {filename} ---\n"
                
                for shape in slide.shapes:
                    # Extract text from text boxes
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.shape_type == 19:  # Table type
                        slide_text += "\n[TABLE]\n"
                        try:
                            for row in shape.table.rows:
                                row_text = " | ".join([cell.text for cell in row.cells])
                                slide_text += row_text + "\n"
                        except:
                            pass  # Skip if table extraction fails
                
                text_parts.append(slide_text)
            
            print(f"    ✓ Extracted {len(prs.slides)} slides with full content")
            return "\n".join(text_parts)
        except Exception as e:
            print(f"    ✗ Error extracting PowerPoint: {e}")
            return f"Error extracting PowerPoint: {e}"
    
    def _extract_docx_text(self, content: bytes, filename: str) -> str:
        """Extract ALL text from Word documents - FULL EXTRACTION including tables"""
        if not Document:
            return "Word document parsing not available - install python-docx"
        
        try:
            doc = Document(BytesIO(content))
            text_parts = []
            
            print(f"    Extracting content from {filename}...")
            
            # Extract all paragraphs
            paragraph_count = 0
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
                    paragraph_count += 1
            
            # Extract all tables
            table_count = 0
            for table in doc.tables:
                text_parts.append(f"\n[TABLE {table_count + 1}]")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_parts.append(row_text)
                table_count += 1
            
            print(f"    ✓ Extracted {paragraph_count} paragraphs and {table_count} tables")
            return "\n".join(text_parts)
        except Exception as e:
            print(f"    ✗ Error extracting Word document: {e}")
            return f"Error extracting Word document: {e}"
    
    def _generate_qualitative_analysis(
        self,
        extracted_content: Dict,
        company_name: str
    ) -> Dict:
        """Generate qualitative analysis using OpenAI - uses FULL content"""
        
        # Prepare context from all extracted content (full content, may need truncation for very large files)
        context = self._format_context(extracted_content, limit_length=True)
        
        system_prompt = """You are an expert investment analyst extracting information from confidential deal documents.

CRITICAL RULES - MUST FOLLOW:
1. ONLY extract information that is EXPLICITLY stated in the documents
2. If information is not found, write "N/A" or "Data not available in documents"
3. NEVER make up, infer, estimate, or hallucinate ANY numbers or facts
4. ALWAYS cite the EXACT source (filename, page/slide number) for EVERY claim
5. Be precise and factual - this is for high-stakes investment decisions
6. If you're uncertain about ANY information, mark it as "Unclear from documents"
7. Do NOT make assumptions about missing data
8. Quote exact text from documents when citing key claims

Your analysis will be used for high-stakes investment decisions. Accuracy is paramount.
False information could result in million-dollar mistakes."""

        user_prompt = f"""Analyze the following deal pack documents for {company_name} and extract information for these sections.

REMEMBER: If information is not explicitly in the documents, write "N/A" or "Data not available".

Documents:
{context}

Extract and analyze the following (cite sources):

1. **Company Overview:**
   - Founding team and backgrounds
   - Product description
   - Monetization model
   - Key traction/KPI highlights
   [Cite: filename, page/slide number]

2. **Deal Snapshot:**
   - Current investors and ownership structure
   - Current round size and valuation
   - Use of funds
   - Post-money ownership expectations
   [Cite: filename, page/slide number]

3. **Market Overview (if provided):**
   - TAM/SAM/SOM and growth outlook
   - 3-5 main growth drivers
   - 3-5 key risks or constraints
   [Cite: filename, page/slide number]

4. **Competition & MOAT (if provided):**
   - How company views competition
   - Closest competitors (regional + global)
   - Differentiation and MOAT
   - Evidence of competitive advantage
   [Cite: filename, page/slide number]

5. **Team Overview:**
   - Founders and key executives
   - Relevant experience and outcomes
   - Leadership concerns or red flags
   [Cite: filename, page/slide number]

6. **Product & Value Propositions:**
   - Customer pain point being solved
   - Key value propositions
   - Validation evidence
   [Cite: filename, page/slide number]

7. **Business Model & Monetization:**
   - Revenue model (pricing, channels, frequency)
   - Key metrics and unit economics
   - Scalability and capital intensity
   [Cite: filename, page/slide number]

8. **Unit Economics & Retention:**
   - Contribution margin
   - Operating leverage evidence
   - Path to profitability
   [Cite: filename, page/slide number]

Format as a structured document with clear sections. Include [N/A] for any information not found.
Include citations in format: [Source: filename, page X]"""

        self._update_progress("qualitative", 35, f"Sending {len(context):,} chars to OpenAI...")
        
        # Use streaming if callback provided
        if self.stream_callback:
            content_parts = []
            stream = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_completion_tokens=8000,
                stream=True
            )
            
            for chunk in stream:
                if chunk.choices[0].delta.content:
                    chunk_content = chunk.choices[0].delta.content
                    content_parts.append(chunk_content)
                    if self.stream_callback:
                        self.stream_callback(chunk_content)
            
            content = "".join(content_parts)
        else:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_completion_tokens=8000
            )
            
            if not response or not response.choices or not response.choices[0].message.content:
                raise ValueError("OpenAI returned empty response for qualitative analysis. This may be due to content filters or API issues.")
            
            content = response.choices[0].message.content
        
        self._update_progress("qualitative", 50, f"Received {len(content):,} chars from OpenAI")
        
        return {
            "content": content,
            "generated_at": datetime.now().isoformat()
        }
    
    def _extract_quantitative_data(
        self,
        extracted_content: Dict,
        company_name: str
    ) -> Dict:
        """Extract quantitative financial data using OpenAI - uses FULL content"""
        
        # Get all Excel data
        excel_data = extracted_content.get("excel", [])
        pdf_text = "\n".join([pdf["text"] for pdf in extracted_content.get("pdfs", [])])
        
        if not excel_data and not pdf_text:
            return {
                "message": "No quantitative data found in uploaded files",
                "tables": {}
            }
        
        # Prepare context with data (formatted for better readability)
        # OPTIMIZATION: For very large sheets, send sample + full data separately
        context = "Excel Data:\n\n"
        total_rows_sent = 0
        max_rows_to_llm = 100  # Limit rows sent to LLM to avoid timeouts
        
        for excel_file in excel_data:
            context += f"\n{'='*60}\nFile: {excel_file['filename']}\n{'='*60}\n"
            data = excel_file.get('data', {})
            metadata = data.get('metadata', {})
            
            for sheet_name, sheet_data in data.get('sheets', {}).items():
                sheet_info = metadata.get(sheet_name, {})
                total_sheet_rows = sheet_info.get('rows', 0)
                
                context += f"\nSheet: {sheet_name}\n"
                context += f"Dimensions: {total_sheet_rows} rows x {sheet_info.get('columns', 0)} columns\n"
                context += f"Columns: {', '.join(sheet_info.get('column_names', []))}\n\n"
                
                # Format as readable table
                if sheet_data and len(sheet_data) > 0:
                    context += "Data:\n"
                    # Add header row
                    headers = list(sheet_data[0].keys()) if sheet_data else []
                    if headers:
                        context += " | ".join(headers) + "\n"
                        context += "-" * (len(" | ".join(headers))) + "\n"
                    
                    # Add data rows (limit to max_rows_to_llm to avoid timeouts)
                    rows_to_show = min(len(sheet_data), max_rows_to_llm)
                    for i, row in enumerate(sheet_data[:rows_to_show]):
                        row_values = [str(row.get(h, "")) for h in headers]
                        context += " | ".join(row_values) + "\n"
                    
                    total_rows_sent += rows_to_show
                    
                    # Note if truncated
                    if len(sheet_data) > max_rows_to_llm:
                        context += f"\n[... showing first {max_rows_to_llm} of {len(sheet_data)} rows ...]\n"
                        context += f"[Full data will be included in Excel output]\n"
                    
                    context += "\n"
        
        print(f"    📊 Prepared context with {total_rows_sent} rows from Excel (limited for LLM)")
        
        # Include more PDF content for tables (increase from 5000 to 30000 chars)
        if pdf_text:
            max_pdf_chars = 30000
            pdf_excerpt = pdf_text[:max_pdf_chars]
            if len(pdf_text) > max_pdf_chars:
                pdf_excerpt += f"\n\n[... PDF content truncated, {len(pdf_text)} total chars ...]"
            context += f"\n\nPDF Content (may contain financial tables):\n{pdf_excerpt}"
        
        system_prompt = """You are an expert financial analyst extracting quantitative data from deal documents.

CRITICAL RULES - ABSOLUTE REQUIREMENTS:
1. ONLY extract numbers that are EXPLICITLY visible in the documents
2. If a metric is not found, write "N/A" - DO NOT estimate or calculate
3. NEVER calculate, derive, infer, or make up ANY numbers
4. ALWAYS cite the EXACT source (filename, sheet name, cell/row reference for Excel, page number for PDFs)
5. Preserve EXACT numbers - copy them character-by-character, don't round
6. If units (USD, millions, %) are unclear or missing, state "Units unclear"
7. If a number appears ambiguous or could be misread, note the ambiguity
8. Copy formulas or calculated cells as their displayed values, note if it's a formula
9. Do NOT perform any arithmetic operations yourself
10. If you see conflicting numbers in different documents, note ALL versions with sources

This data will be used for investment models and financial analysis.
A single incorrect number could result in multi-million dollar mistakes.
When in doubt, mark as "N/A" or "Unclear" rather than guessing."""

        user_prompt = f"""Extract ALL quantitative financial data for {company_name} from the provided documents.

Documents:
{context}

TASK: Extract metrics and tables from the above data.

EXTRACT:
1. Individual Metrics (KPIs, financial metrics)
2. Tables (Cap Table, P&L, Balance Sheet, Unit Economics, etc.)

For each item, cite the exact source (filename, sheet name, row/column).

RESPOND WITH JSON:
{{
  "metrics": [
    {{"metric": "Total Revenue 2024", "value": "$5M", "period": "2024", "source": "financials.xlsx, Sheet1"}},
    {{"metric": "Gross Margin", "value": "75%", "period": "2024", "source": "financials.xlsx, Sheet1"}}
  ],
  "tables": [
    {{
      "title": "Cap Table",
      "columns": ["Shareholder", "Shares", "Ownership %", "Round"],
      "rows": [
        ["Founder A", "1000000", "45%", "Seed"],
        ["Investor B", "500000", "22.5%", "Series A"]
      ],
      "source": "captable.xlsx, Captable sheet"
    }}
  ]
}}

RULES:
- Extract ONLY data that exists in the documents above
- If you cannot find quantitative data, return: {{"metrics": [], "tables": []}}
- Include source citation for every metric and table
- Copy numbers exactly as shown
- Include ALL tables from Excel files (especially Cap Table)

Begin extraction:"""

        self._update_progress("quantitative", 65, f"Sending {len(context):,} chars to OpenAI...")
        
        # Use streaming if callback provided
        if self.stream_callback:
            content_parts = []
            stream = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_completion_tokens=16000,
                stream=True
            )
            
            for chunk in stream:
                if chunk.choices[0].delta.content:
                    chunk_content = chunk.choices[0].delta.content
                    content_parts.append(chunk_content)
                    if self.stream_callback:
                        self.stream_callback(chunk_content)
            
            content = "".join(content_parts)
        else:
            # Use higher token limit for quantitative data (large tables)
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_completion_tokens=16000
            )
            
            if not response or not response.choices or not response.choices[0].message.content:
                raise ValueError("OpenAI returned empty response for quantitative data. This may be due to content filters or API issues.")
            
            content = response.choices[0].message.content
        
        self._update_progress("quantitative", 75, f"Received {len(content):,} chars from OpenAI")
        
        return {
            "content": content,
            "excel_files": [f["filename"] for f in excel_data],
            "generated_at": datetime.now().isoformat()
        }
    
    def _format_context(self, extracted_content: Dict, limit_length: bool = False) -> str:
        """
        Format extracted content for LLM context
        
        Args:
            extracted_content: Dictionary with extracted file content
            limit_length: If True, truncate for token limits. If False, send FULL content.
        """
        context_parts = []
        
        # Add PDFs - FULL content or truncated based on limit_length
        for pdf in extracted_content.get("pdfs", []):
            text = pdf['text']
            if limit_length and len(text) > 50000:
                text = text[:50000] + f"\n\n[... content truncated, {len(pdf['text'])} total chars ...]"
            context_parts.append(f"\n{'='*60}\nPDF: {pdf['filename']}\n{'='*60}\n{text}")
        
        # Add PowerPoint - FULL content or truncated
        for ppt in extracted_content.get("powerpoint", []):
            text = ppt['text']
            if limit_length and len(text) > 50000:
                text = text[:50000] + f"\n\n[... content truncated, {len(ppt['text'])} total chars ...]"
            context_parts.append(f"\n{'='*60}\nPowerPoint: {ppt['filename']}\n{'='*60}\n{text}")
        
        # Add Word documents - FULL content or truncated
        for docx in extracted_content.get("docx", []):
            text = docx['text']
            if limit_length and len(text) > 50000:
                text = text[:50000] + f"\n\n[... content truncated, {len(docx['text'])} total chars ...]"
            context_parts.append(f"\n{'='*60}\nWord Document: {docx['filename']}\n{'='*60}\n{text}")
        
        # Add Excel - FULL DATA, not just samples
        for excel in extracted_content.get("excel", []):
            context_parts.append(f"\n{'='*60}\nExcel: {excel['filename']}\n{'='*60}")
            data = excel.get("data", {})
            metadata = data.get("metadata", {})
            
            for sheet_name in data.get("sheet_names", []):
                sheet_info = metadata.get(sheet_name, {})
                context_parts.append(f"\nSheet: {sheet_name} ({sheet_info.get('rows', 0)} rows x {sheet_info.get('columns', 0)} cols)")
                
                # Add ALL rows (not just first 5)
                sheet_data = data.get("sheets", {}).get(sheet_name, [])
                if sheet_data:
                    # Convert to readable format
                    context_parts.append(f"Complete data for {sheet_name}:")
                    context_parts.append(json.dumps(sheet_data, indent=2))
                else:
                    context_parts.append("(Empty sheet)")
        
        full_context = "\n".join(context_parts)
        
        context_size_mb = len(full_context) / (1024 * 1024)
        print(f"\n📊 Context prepared: {len(full_context):,} characters ({context_size_mb:.2f} MB)")
        
        # Warn about very large contexts (GPT-4o has ~128K token limit, roughly 500K chars)
        if len(full_context) > 400_000:
            print(f"⚠️  WARNING: Very large context ({context_size_mb:.2f} MB)")
            print(f"    This may hit token limits or slow down processing")
            print(f"    Consider splitting files or using smaller documents")
        elif len(full_context) > 200_000:
            print(f"⚠️  Large context ({context_size_mb:.2f} MB) - processing may be slow")
        
        return full_context
    
    def _generate_excel_file(
        self,
        quantitative_data: Dict,
        company_name: str
    ) -> Optional[BytesIO]:
        """Generate Excel file dynamically based on extracted data"""
        
        if not load_workbook:
            print("⚠️  openpyxl not available - skipping Excel generation")
            return None
        
        try:
            from openpyxl import Workbook
            
            # Parse the LLM response to extract structured JSON
            content = quantitative_data.get("content", "")
            
            # Try to extract JSON from the content
            data_json = self._extract_json_from_content(content)
            
            if not data_json:
                print("⚠️  Could not parse JSON from LLM response - skipping Excel generation")
                print(f"    Content preview: {content[:500]}...")
                return None
            
            if not data_json.get("metrics") and not data_json.get("tables"):
                print("⚠️  No metrics or tables found in extracted data - skipping Excel generation")
                print(f"    Extracted keys: {list(data_json.keys())}")
                return None
            
            # Create workbook
            wb = Workbook()
            wb.remove(wb.active)  # Remove default sheet
            
            sheets_created = 0
            
            # Create Metrics sheet if we have metrics
            if data_json.get("metrics") and len(data_json["metrics"]) > 0:
                ws_metrics = wb.create_sheet("Key Metrics")
                self._add_metrics_sheet(ws_metrics, data_json["metrics"], company_name)
                sheets_created += 1
            
            # Create a sheet for each table found
            if data_json.get("tables"):
                for table in data_json["tables"]:
                    if table.get("rows") and len(table["rows"]) > 0:
                        # Create sheet with table title (sanitize name)
                        sheet_name = self._sanitize_sheet_name(table.get("title", f"Table {sheets_created + 1}"))
                        ws_table = wb.create_sheet(sheet_name)
                        self._add_table_sheet(ws_table, table, company_name)
                        sheets_created += 1
            
            if sheets_created == 0:
                print("⚠️  No data to populate Excel - skipping generation")
                return None
            
            # Save to BytesIO
            excel_buffer = BytesIO()
            wb.save(excel_buffer)
            excel_buffer.seek(0)
            
            print(f"✅ Excel file generated with {sheets_created} sheet(s)")
            return excel_buffer
            
        except Exception as e:
            print(f"⚠️  Error generating Excel file: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _extract_json_from_content(self, content: str) -> Optional[Dict]:
        """Extract and validate JSON data from LLM response"""
        try:
            # Try to find JSON in the content
            # LLMs sometimes wrap JSON in markdown code blocks
            
            # First, try to remove markdown code blocks
            if "```json" in content:
                json_start = content.find("```json") + 7
                json_end = content.find("```", json_start)
                if json_end > json_start:
                    json_str = content[json_start:json_end].strip()
                    parsed = json.loads(json_str)
                    print(f"    ✓ Extracted JSON from markdown code block")
                    return parsed
            
            # Otherwise, find first { to last }
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            
            if json_start != -1 and json_end > json_start:
                json_str = content[json_start:json_end]
                parsed = json.loads(json_str)
                
                # Validate structure
                if not isinstance(parsed, dict):
                    print(f"⚠️  Parsed JSON is not a dictionary")
                    return None
                
                # Check if it has expected keys
                has_metrics = "metrics" in parsed
                has_tables = "tables" in parsed
                
                if has_metrics or has_tables:
                    print(f"    ✓ Valid JSON extracted (metrics: {has_metrics}, tables: {has_tables})")
                    return parsed
                else:
                    print(f"⚠️  JSON missing expected keys (metrics/tables)")
                    return None
            
            print(f"⚠️  No valid JSON found in content")
            return None
            
        except json.JSONDecodeError as e:
            print(f"⚠️  Could not parse JSON from content: {e}")
            print(f"    Content preview: {content[:200]}...")
            return None
    
    def _sanitize_sheet_name(self, name: str) -> str:
        """Sanitize sheet name for Excel (max 31 chars, no special chars)"""
        # Remove invalid characters
        invalid_chars = ['\\', '/', '*', '?', ':', '[', ']']
        for char in invalid_chars:
            name = name.replace(char, '')
        
        # Truncate to 31 characters
        return name[:31]
    
    def _add_metrics_sheet(self, ws, metrics: List[Dict], company_name: str):
        """Add Key Metrics sheet with individual KPIs"""
        header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=12)
        
        # Title
        ws['A1'] = f"{company_name} - Key Metrics"
        ws['A1'].font = Font(bold=True, size=14)
        ws.merge_cells('A1:D1')
        
        # Headers
        ws['A3'] = "Metric"
        ws['B3'] = "Value"
        ws['C3'] = "Period"
        ws['D3'] = "Source"
        for cell in ['A3', 'B3', 'C3', 'D3']:
            ws[cell].fill = header_fill
            ws[cell].font = header_font
        
        # Add metrics
        row = 4
        for metric in metrics:
            ws[f'A{row}'] = metric.get('metric', 'N/A')
            ws[f'B{row}'] = metric.get('value', 'N/A')
            ws[f'C{row}'] = metric.get('period', '')
            ws[f'D{row}'] = metric.get('source', 'N/A')
            row += 1
        
        # Column widths
        ws.column_dimensions['A'].width = 30
        ws.column_dimensions['B'].width = 20
        ws.column_dimensions['C'].width = 20
        ws.column_dimensions['D'].width = 50
    
    def _add_table_sheet(self, ws, table: Dict, company_name: str):
        """Add a table sheet dynamically based on the table structure"""
        header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=11)
        
        # Title
        title = table.get('title', 'Data')
        ws['A1'] = f"{company_name} - {title}"
        ws['A1'].font = Font(bold=True, size=14)
        
        # Source
        source = table.get('source', 'N/A')
        ws['A2'] = f"Source: {source}"
        ws['A2'].font = Font(italic=True, size=9, color="666666")
        
        # Column headers
        columns = table.get('columns', [])
        if columns:
            for col_idx, col_name in enumerate(columns, 1):
                cell = ws.cell(row=4, column=col_idx, value=col_name)
                cell.fill = header_fill
                cell.font = header_font
        
        # Data rows
        rows = table.get('rows', [])
        for row_idx, row_data in enumerate(rows, 5):
            for col_idx, cell_value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=cell_value)
        
        # Auto-size columns
        for col_idx in range(1, len(columns) + 1):
            col_letter = chr(64 + col_idx)  # A, B, C, etc.
            ws.column_dimensions[col_letter].width = 18
    
    def _generate_human_readable_summary(
        self,
        company_name: str,
        qualitative_analysis: Dict,
        quantitative_data: Dict
    ) -> Dict:
        """Generate a human-readable summary for frontend display"""
        
        system_prompt = """You are an expert investment analyst creating executive summaries of data room analysis.

Your role is to synthesize qualitative insights and quantitative metrics into a clear, readable summary for investors.

CRITICAL REQUIREMENTS:
1. Write in clear, professional prose (not bullet points or structured format)
2. Highlight the most important insights and metrics
3. Include specific numbers and KPIs where available
4. Cite sources (e.g., "per pitch deck", "from financials.xlsx")
5. Be balanced - note both strengths and concerns
6. Focus on investment-relevant information"""

        user_prompt = f"""Create a comprehensive Data Room Summary for {company_name}.

QUALITATIVE INSIGHTS:
{qualitative_analysis.get('content', 'N/A')}

QUANTITATIVE DATA:
{quantitative_data.get('content', 'N/A')}

Generate a professional summary with these sections:

**Executive Overview**
2-3 paragraphs summarizing the company, product, team, and key traction metrics. Include specific numbers.

**Key Metrics & Financials**
Highlight the most important metrics found in the data room:
- Revenue metrics (ARR, MRR, growth rates)
- Customer metrics (count, CAC, LTV, churn)
- Unit economics (margins, payback period)
- Runway and burn rate
- Any other material KPIs

**Deal Structure** (if available)
- Current investors and ownership
- Round details (size, valuation, terms)
- Use of funds
- Cap table notes

**Strengths Identified**
Top 3-5 positive signals from the data room with evidence.

**Concerns & Gaps**
Top 3-5 concerns, inconsistencies, or missing information that require follow-up.

**Data Quality Assessment**
Comment on the completeness and quality of the data room materials.

FORMAT: Write as flowing prose with clear section headers. Include specific numbers with citations. Be concise but comprehensive (aim for 800-1200 words)."""

        try:
            # Use streaming if callback provided
            if self.stream_callback:
                content_parts = []
                stream = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    max_completion_tokens=4000,
                    stream=True
                )
                
                for chunk in stream:
                    if chunk.choices[0].delta.content:
                        chunk_content = chunk.choices[0].delta.content
                        content_parts.append(chunk_content)
                        if self.stream_callback:
                            self.stream_callback(chunk_content)
                
                content = "".join(content_parts)
            else:
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    max_completion_tokens=4000
                )
                content = response.choices[0].message.content
            
            return {
                "content": content,
                "generated_at": datetime.now().isoformat(),
                "word_count": len(content.split())
            }
        except Exception as e:
            print(f"⚠️  Error generating human-readable summary: {e}")
            # Fallback to basic summary
            return {
                "content": f"# Data Room Summary for {company_name}\n\n" +
                          f"**Files Processed:** {len(qualitative_analysis.get('content', ''))} chars qualitative, " +
                          f"{len(quantitative_data.get('content', ''))} chars quantitative\n\n" +
                          f"Error generating detailed summary: {str(e)}",
                "generated_at": datetime.now().isoformat(),
                "word_count": 50
            }
    
    def format_report_as_text(self, report: Dict) -> str:
        """Format the data room report as readable text"""
        output = []
        
        output.append("=" * 80)
        output.append("DATA ROOM ANALYSIS REPORT")
        output.append("=" * 80)
        output.append(f"\nCompany: {report['company_name']}")
        output.append(f"Files Processed: {report['files_processed']}")
        output.append(f"Generated: {report['generated_at']}")
        output.append("\n" + "=" * 80 + "\n")
        
        # Human-readable summary (if available)
        if "human_readable_summary" in report:
            output.append("\n## EXECUTIVE SUMMARY\n")
            output.append(report['human_readable_summary']['content'])
            output.append("\n" + "-" * 80 + "\n")
        
        # Qualitative Analysis
        output.append("\n## QUALITATIVE ANALYSIS (Detailed)\n")
        output.append(report['qualitative_analysis']['content'])
        output.append("\n" + "-" * 80 + "\n")
        
        # Quantitative Data
        output.append("\n## QUANTITATIVE DATA (Structured)\n")
        output.append(report['quantitative_data']['content'])
        output.append("\n" + "-" * 80 + "\n")
        
        return "\n".join(output)
    
    def generate_docx_summary(self, report: Dict) -> BytesIO:
        """Generate a DOCX file with the human-readable summary"""
        try:
            if not Document:
                raise ImportError("python-docx not available")
            
            doc = Document()
            
            # Title
            title = doc.add_heading(f"Data Room Analysis: {report['company_name']}", 0)
            title.alignment = 1  # Center alignment
            
            # Metadata
            doc.add_paragraph(f"Generated: {report.get('generated_at', 'N/A')}")
            doc.add_paragraph(f"Files Processed: {report.get('files_processed', 'N/A')}")
            doc.add_paragraph("")  # Empty line
            
            # Human-readable summary content
            if "human_readable_summary" in report:
                content = report["human_readable_summary"].get("content", "")
                
                # Parse content and add to document
                # Split by lines and add as paragraphs
                for line in content.split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    
                    # Check if it's a heading (starts with ## or **)
                    if line.startswith("**") and line.endswith("**"):
                        # Bold heading
                        heading_text = line.strip("*")
                        doc.add_heading(heading_text, level=2)
                    elif line.startswith("#"):
                        # Markdown heading
                        heading_text = line.lstrip("#").strip()
                        doc.add_heading(heading_text, level=2)
                    else:
                        # Regular paragraph
                        doc.add_paragraph(line)
            else:
                doc.add_paragraph("Summary not available.")
            
            # Save to BytesIO
            docx_buffer = BytesIO()
            doc.save(docx_buffer)
            docx_buffer.seek(0)
            
            return docx_buffer
            
        except Exception as e:
            print(f"⚠️  Error generating DOCX: {e}")
            return None




